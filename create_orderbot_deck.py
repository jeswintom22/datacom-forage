from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

OUTPUT = "OrderBot_Presentation_refined.pptx"

NAVY = RGBColor(18, 32, 63)
BLUE = RGBColor(46, 92, 170)
TEAL = RGBColor(28, 147, 137)
GREEN = RGBColor(47, 142, 79)
RED = RGBColor(196, 61, 61)
AMBER = RGBColor(200, 132, 29)
GRAY = RGBColor(98, 106, 120)
LIGHT = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
DARK = RGBColor(33, 37, 41)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

LEFT = Inches(0.6)
RIGHT = Inches(12.73)
CONTENT_W = Inches(12.13)


def add_bg(slide, accent):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = WHITE
    bg.line.fill.background()

    accent_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.14), prs.slide_height)
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = accent
    accent_bar.line.fill.background()

    top = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.14), 0, prs.slide_width - Inches(0.14), Inches(0.42))
    top.fill.solid()
    top.fill.fore_color.rgb = accent
    top.line.fill.background()

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.9), prs.slide_width, Inches(0.6))
    band.fill.solid()
    band.fill.fore_color.rgb = LIGHT
    band.line.fill.background()


def add_title(slide, title, subtitle=None, accent=BLUE):
    title_box = slide.shapes.add_textbox(LEFT, Inches(0.7), Inches(9.0), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.size = Pt(26)
    r.font.bold = True
    r.font.color.rgb = NAVY
    r.font.name = "Aptos Display"
    if subtitle:
        sub = slide.shapes.add_textbox(LEFT, Inches(1.34), Inches(9.8), Inches(0.5))
        tf2 = sub.text_frame
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.size = Pt(12.5)
        r2.font.color.rgb = GRAY
        r2.font.name = "Aptos"


def add_textbox(slide, left, top, width, height, text, size=16, color=DARK, bold=False, bullet=False, fill=None, line=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    if fill is not None:
        box.fill.solid()
        box.fill.fore_color.rgb = fill
    else:
        box.fill.background()
    if line is not None:
        box.line.color.rgb = line
    else:
        box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    if isinstance(text, list):
        for i, item in enumerate(text):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = item
            p.level = 0
            p.bullet = bullet
            p.font.size = Pt(size)
            p.font.color.rgb = color
            p.font.name = "Aptos"
            p.font.bold = bold
            p.alignment = align
    else:
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.name = "Aptos"
        p.font.bold = bold
        p.alignment = align
    return box


def add_card(slide, left, top, width, height, title, body, fill=LIGHT, title_color=NAVY, body_color=DARK, title_size=15, body_size=12, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line if line is not None else fill
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True

    p1 = tf.paragraphs[0]
    p1.space_after = Pt(4)
    r1 = p1.add_run()
    r1.text = title
    r1.font.size = Pt(title_size)
    r1.font.bold = True
    r1.font.color.rgb = title_color
    r1.font.name = "Aptos"

    if body:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(0)
        r2 = p2.add_run()
        r2.text = body
        r2.font.size = Pt(body_size)
        r2.font.color.rgb = body_color
        r2.font.name = "Aptos"
    return shape


def add_box(slide, left, top, width, height, text, fill, font_color=WHITE, line=None, size=14, bold=True, rounded=True):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.color.rgb = fill
    else:
        shape.line.color.rgb = line
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = font_color
    r.font.name = "Aptos"
    return shape


def add_footer(slide, accent_label):
    footer_left = slide.shapes.add_textbox(LEFT, Inches(7.03), Inches(4.0), Inches(0.35))
    p = footer_left.text_frame.paragraphs[0]
    r = p.add_run()
    r.text = accent_label
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = GRAY
    r.font.name = "Aptos"

    footer_right = slide.shapes.add_textbox(Inches(11.6), Inches(7.01), Inches(1.0), Inches(0.25))
    p2 = footer_right.text_frame.paragraphs[0]
    p2.alignment = PP_ALIGN.RIGHT
    r2 = p2.add_run()
    r2.text = "OrderBot"
    r2.font.size = Pt(10)
    r2.font.bold = True
    r2.font.color.rgb = GRAY
    r2.font.name = "Aptos"


def arrow(slide, x1, y1, x2, y2, color=GRAY):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(2)
    return conn


def add_badge(slide, left, top, text, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, Inches(1.0), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.size = Pt(18)
    r.font.bold = True
    r.font.color.rgb = WHITE
    r.font.name = "Aptos"
    return shape


# Slide 1: problem
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, RED)
add_title(slide, "Slide 1. The Problem: The Current Order Process Is a Bottleneck",
          "Manual handoffs create delays, errors, and a dependency on staff availability.")

add_card(slide, LEFT, Inches(1.75), Inches(12.05), Inches(0.82), "Manual workflow", "Orders move through email, people, Salesforce, and spreadsheets before a response is sent.", fill=LIGHT, title_size=16, body_size=12, line=RED)

add_box(slide, Inches(0.72), Inches(2.95), Inches(1.42), Inches(0.72), "Email\nPDF arrives", RED, size=14)
add_box(slide, Inches(2.35), Inches(2.95), Inches(1.42), Inches(0.72), "Manual\nreading", AMBER, size=14)
add_box(slide, Inches(3.98), Inches(2.95), Inches(1.42), Inches(0.72), "Salesforce\ncheck", BLUE, size=14)
add_box(slide, Inches(5.61), Inches(2.95), Inches(1.58), Inches(0.72), "Inventory\ncheck", BLUE, size=14)
add_box(slide, Inches(7.43), Inches(2.95), Inches(1.38), Inches(0.72), "Decision", NAVY, size=15)
add_box(slide, Inches(9.05), Inches(2.65), Inches(1.1), Inches(0.56), "OK", GREEN, size=15)
add_box(slide, Inches(9.05), Inches(3.37), Inches(1.1), Inches(0.56), "Issue", RED, size=15)
add_box(slide, Inches(10.35), Inches(2.58), Inches(2.3), Inches(0.68), "Send customer and warehouse emails", TEAL, size=13)
add_box(slide, Inches(10.35), Inches(3.32), Inches(2.3), Inches(0.68), "Email sales rep and wait", AMBER, size=13)

arrow(slide, Inches(2.14), Inches(3.31), Inches(2.35), Inches(3.31))
arrow(slide, Inches(3.78), Inches(3.31), Inches(3.98), Inches(3.31))
arrow(slide, Inches(5.41), Inches(3.31), Inches(5.61), Inches(3.31))
arrow(slide, Inches(7.19), Inches(3.31), Inches(7.43), Inches(3.31))
arrow(slide, Inches(8.81), Inches(3.1), Inches(9.05), Inches(2.93))
arrow(slide, Inches(8.81), Inches(3.48), Inches(9.05), Inches(3.65))
arrow(slide, Inches(10.15), Inches(2.93), Inches(10.35), Inches(2.93))
arrow(slide, Inches(10.15), Inches(3.65), Inches(10.35), Inches(3.65))

add_card(slide, Inches(0.72), Inches(4.15), Inches(3.82), Inches(1.35), "Manual effort", "Each order is re-read and re-keyed by hand before action can be taken.", fill=LIGHT, title_size=15, body_size=12, line=RED)
add_card(slide, Inches(4.75), Inches(4.15), Inches(3.82), Inches(1.35), "Delay risk", "If staff are busy or away, orders sit in the inbox and customers wait longer.", fill=LIGHT, title_size=15, body_size=12, line=AMBER)
add_card(slide, Inches(8.78), Inches(4.15), Inches(3.77), Inches(1.35), "Error risk", "Information is copied across systems, which raises the chance of missed details.", fill=LIGHT, title_size=15, body_size=12, line=BLUE)

add_badge(slide, Inches(10.95), Inches(5.62), "HOURS", RED)
add_textbox(slide, Inches(10.45), Inches(6.35), Inches(2.4), Inches(0.3), "Typical processing time today", size=11.5, color=GRAY, align=PP_ALIGN.CENTER)
add_footer(slide, "Current state")

# Slide 2: solution
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, TEAL)
add_title(slide, "Slide 2. The Solution: OrderBot Automates the Workflow",
          "A goal-driven agent reads the order, checks systems, makes the routine decision, and routes exceptions to humans.")

add_card(slide, Inches(0.75), Inches(1.82), Inches(3.2), Inches(1.62), "1. Detect & read", "Monitors the shared inbox, opens the PDF, and extracts the order details automatically.", fill=LIGHT, title_size=16, body_size=12, line=TEAL)
add_card(slide, Inches(4.2), Inches(1.82), Inches(3.2), Inches(1.62), "2. Verify & compare", "Checks the account in Salesforce and confirms stock levels in the inventory sheet.", fill=LIGHT, title_size=16, body_size=12, line=BLUE)
add_card(slide, Inches(7.65), Inches(1.82), Inches(3.2), Inches(1.62), "3. Decide & act", "Sends confirmations when everything matches. Routes exceptions to a human reviewer with context.", fill=LIGHT, title_size=16, body_size=12, line=GREEN)
add_card(slide, Inches(11.1), Inches(1.82), Inches(1.5), Inches(1.62), "24/7", "Always on", fill=TEAL, title_size=18, body_size=12, line=TEAL)

add_box(slide, Inches(0.9), Inches(4.35), Inches(1.55), Inches(0.78), "Inbox", TEAL, size=16)
add_box(slide, Inches(2.7), Inches(4.35), Inches(1.55), Inches(0.78), "Parser", BLUE, size=16)
add_box(slide, Inches(4.5), Inches(4.35), Inches(1.55), Inches(0.78), "Salesforce", BLUE, size=15)
add_box(slide, Inches(6.3), Inches(4.35), Inches(1.75), Inches(0.78), "Inventory", BLUE, size=16)
add_box(slide, Inches(8.3), Inches(4.35), Inches(1.5), Inches(0.78), "Decision", NAVY, size=16)
add_box(slide, Inches(10.1), Inches(4.0), Inches(1.45), Inches(0.62), "OK", GREEN, size=15)
add_box(slide, Inches(10.1), Inches(4.82), Inches(1.45), Inches(0.62), "Exception", AMBER, size=14)
add_box(slide, Inches(11.8), Inches(4.0), Inches(1.05), Inches(0.62), "Send", TEAL, size=15)
add_box(slide, Inches(11.8), Inches(4.82), Inches(1.05), Inches(0.62), "Review", RED, size=14)

arrow(slide, Inches(2.45), Inches(4.74), Inches(2.7), Inches(4.74))
arrow(slide, Inches(4.25), Inches(4.74), Inches(4.5), Inches(4.74))
arrow(slide, Inches(6.05), Inches(4.74), Inches(6.3), Inches(4.74))
arrow(slide, Inches(7.95), Inches(4.74), Inches(8.3), Inches(4.74))
arrow(slide, Inches(9.8), Inches(4.5), Inches(10.1), Inches(4.31))
arrow(slide, Inches(9.8), Inches(4.87), Inches(10.1), Inches(5.13))
arrow(slide, Inches(11.55), Inches(4.31), Inches(11.8), Inches(4.31))
arrow(slide, Inches(11.55), Inches(5.13), Inches(11.8), Inches(5.13))

add_card(slide, Inches(0.85), Inches(5.55), Inches(12.0), Inches(0.82), "OrderBot in plain English", "It reads the order, checks the systems, and either completes the routine work or escalates only the exceptions that need people.", fill=LIGHT, title_size=15, body_size=12, line=TEAL)
add_footer(slide, "Target outcome")

# Slide 3: impact
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, NAVY)
add_title(slide, "Slide 3. Business Impact & Recommended Next Steps",
          "Automation reduces cycle time, improves consistency, and gives the team back time for higher-value work.")

add_card(slide, Inches(0.8), Inches(1.82), Inches(3.0), Inches(1.42), "Processing time", "Hours → Minutes", fill=LIGHT, title_size=16, body_size=20, line=BLUE)
add_card(slide, Inches(4.0), Inches(1.82), Inches(3.0), Inches(1.42), "Accuracy", "Fewer manual errors", fill=LIGHT, title_size=16, body_size=20, line=TEAL)
add_card(slide, Inches(7.2), Inches(1.82), Inches(3.0), Inches(1.42), "Capacity", "Free ops time", fill=LIGHT, title_size=16, body_size=20, line=GREEN)
add_card(slide, Inches(10.4), Inches(1.82), Inches(2.05), Inches(1.42), "24/7", "Coverage", fill=LIGHT, title_size=18, body_size=20, line=AMBER)

add_card(slide, Inches(0.85), Inches(3.75), Inches(6.0), Inches(1.85), "Expected business benefits", "• Faster confirmation and fulfillment for customers\n• Consistent decisions based on the same rules every time\n• Reduced dependence on individual staff availability\n• More time for exceptions, customer service, and higher-value work", fill=LIGHT, title_size=17, body_size=12.5, line=WHITE)

add_card(slide, Inches(7.15), Inches(3.75), Inches(5.2), Inches(1.85), "Recommended next step", "Run a read-only proof of concept that watches the inbox, parses orders, and compares results against Salesforce and inventory without sending any customer emails.", fill=WHITE, title_color=NAVY, body_color=DARK, title_size=17, body_size=13, line=TEAL)

add_card(slide, Inches(7.15), Inches(5.8), Inches(5.2), Inches(0.66), "Success metrics", "< 5 minute processing • 99.5% accuracy target • clear exception path", fill=TEAL, title_color=WHITE, body_color=WHITE, title_size=14, body_size=12, line=TEAL)
add_footer(slide, "Business impact")

prs.save(OUTPUT)
print(f"Created {OUTPUT}")
